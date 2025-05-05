# app.py
from flask import Flask, request, jsonify, render_template, send_file, session, redirect, url_for
from ollama_client import generate_content
import random
import json
import os
import tempfile
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import textwrap
from datetime import datetime

# Import database models and authentication routes
from models import db, User, Presentation as PresentationModel, Slide
from auth import auth_bp, login_required

app = Flask(__name__)
app.config['SECRET_KEY'] = os.environ.get('SECRET_KEY', 'dev_key_change_in_production')
app.config['SQLALCHEMY_DATABASE_URI'] = os.environ.get('DATABASE_URL', 'sqlite:///pptgenerator.db')
app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

# Initialize database
db.init_app(app)

# Register blueprints
app.register_blueprint(auth_bp, url_prefix='/auth')

# Available layouts
LAYOUTS = [
    "titleAndBullets",
    "quote",
    "imageAndParagraph",
    "twoColumn",
    "titleOnly"
]

# Template definitions
TEMPLATES = {
    'corporate': {
        'colors': {
            'primary': RGBColor(15, 76, 129),    # #0f4c81
            'secondary': RGBColor(110, 156, 196), # #6e9cc4
            'accent': RGBColor(242, 177, 56),    # #f2b138
            'background': RGBColor(255, 255, 255), # #ffffff
            'text': RGBColor(51, 51, 51)         # #333333
        },
        'font': 'Arial'
    },
    'creative': {
        'colors': {
            'primary': RGBColor(255, 107, 107),  # #ff6b6b
            'secondary': RGBColor(78, 205, 196), # #4ecdc4
            'accent': RGBColor(255, 209, 102),   # #ffd166
            'background': RGBColor(249, 241, 230), # #f9f1e6
            'text': RGBColor(90, 57, 33)         # #5a3921
        },
        'font': 'Georgia'
    },
    'minimal': {
        'colors': {
            'primary': RGBColor(44, 62, 80),     # #2c3e50
            'secondary': RGBColor(149, 165, 166), # #95a5a6
            'accent': RGBColor(231, 76, 60),     # #e74c3c
            'background': RGBColor(248, 248, 248), # #f8f8f8
            'text': RGBColor(34, 34, 34)         # #222222
        },
        'font': 'Helvetica'
    },
    'dark': {
        'colors': {
            'primary': RGBColor(187, 134, 252),  # #bb86fc
            'secondary': RGBColor(3, 218, 198),  # #03dac6
            'accent': RGBColor(207, 102, 121),   # #cf6679
            'background': RGBColor(26, 26, 26),  # #1a1a1a
            'text': RGBColor(245, 245, 245)      # #f5f5f5
        },
        'font': 'Roboto'
    }
}

@app.route('/')
def index():
    if 'user_id' in session:
        return redirect(url_for('dashboard'))
    return render_template('index.html')

@app.route('/dashboard')
@login_required
def dashboard():
    user_id = session.get('user_id')
    username = session.get('username')
    
    # Get all presentations for the user
    presentations = PresentationModel.query.filter_by(user_id=user_id).order_by(PresentationModel.updated_at.desc()).all()
    
    return render_template('dashboard.html', 
                          username=username,
                          presentations=presentations)

@app.route('/editor')
@login_required
def editor():
    return render_template('editor.html')

@app.route('/editor/<int:presentation_id>')
@login_required
def edit_presentation(presentation_id):
    user_id = session.get('user_id')
    presentation = PresentationModel.query.filter_by(id=presentation_id, user_id=user_id).first()
    
    if not presentation:
        return redirect(url_for('dashboard'))
    
    return render_template('editor.html', presentation=presentation.to_dict())

@app.route('/api/generate', methods=['POST'])
def generate():
    data = request.json
    template = data.get('template')
    topic = data.get('topic')
    slide_count = data.get('slideCount', 6)  # Default to 6 if not specified
    
    # Limit slide count to reasonable number
    slide_count = min(max(1, slide_count), 10)
    
    # Generate slides
    slides = []
    for _ in range(slide_count):
        # Select a random layout for each slide
        layout = random.choice(LAYOUTS)
        
        # Generate content using Ollama based on the layout and topic
        content = generate_content(layout, topic)
        
        # Process content to prevent overflow
        processed_content = process_content_for_layout(content, layout)
        
        slides.append({
            'layout': layout,
            'content': processed_content
        })
    
    return jsonify({
        'slides': slides,
        'template': template
    })

def process_content_for_layout(content, layout):
    """Process and truncate content based on layout to prevent overflow"""
    processed = dict(content)
    
    if layout == 'titleAndBullets':
        # Limit title length
        if 'title' in processed and len(processed['title']) > 80:
            processed['title'] = processed['title'][:77] + '...'
        
        # Limit number of bullets and length of each
        if 'bullets' in processed:
            # Keep maximum 5 bullets
            processed['bullets'] = processed['bullets'][:5]
            
            # Limit length of each bullet
            processed['bullets'] = [
                (bullet[:97] + '...') if len(bullet) > 100 else bullet 
                for bullet in processed['bullets']
            ]
    
    elif layout == 'quote':
        # Limit quote length
        if 'quote' in processed and len(processed['quote']) > 150:
            processed['quote'] = processed['quote'][:147] + '...'
            
        # Limit author length
        if 'author' in processed and len(processed['author']) > 50:
            processed['author'] = processed['author'][:47] + '...'
    
    elif layout == 'imageAndParagraph':
        # Limit title length
        if 'title' in processed and len(processed['title']) > 80:
            processed['title'] = processed['title'][:77] + '...'
            
        # Limit paragraph length
        if 'paragraph' in processed and len(processed['paragraph']) > 300:
            processed['paragraph'] = processed['paragraph'][:297] + '...'
            
        # Limit image description length
        if 'imageDescription' in processed and len(processed['imageDescription']) > 100:
            processed['imageDescription'] = processed['imageDescription'][:97] + '...'
    
    elif layout == 'twoColumn':
        # Limit title length
        if 'title' in processed and len(processed['title']) > 80:
            processed['title'] = processed['title'][:77] + '...'
            
        # Limit column titles length
        if 'column1Title' in processed and len(processed['column1Title']) > 50:
            processed['column1Title'] = processed['column1Title'][:47] + '...'
        if 'column2Title' in processed and len(processed['column2Title']) > 50:
            processed['column2Title'] = processed['column2Title'][:47] + '...'
            
        # Limit column content length
        if 'column1Content' in processed and len(processed['column1Content']) > 200:
            processed['column1Content'] = processed['column1Content'][:197] + '...'
        if 'column2Content' in processed and len(processed['column2Content']) > 200:
            processed['column2Content'] = processed['column2Content'][:197] + '...'
    
    elif layout == 'titleOnly':
        # Limit title length
        if 'title' in processed and len(processed['title']) > 80:
            processed['title'] = processed['title'][:77] + '...'
            
        # Limit subtitle length
        if 'subtitle' in processed and len(processed['subtitle']) > 120:
            processed['subtitle'] = processed['subtitle'][:117] + '...'
    
    return processed

@app.route('/api/save', methods=['POST'])
@login_required
def save_presentation():
    user_id = session.get('user_id')
    data = request.json
    
    # Extract presentation data
    topic = data.get('topic')
    template_id = data.get('template')
    slides = data.get('slides', [])
    presentation_id = data.get('id')
    
    # Validate input
    if not topic or not template_id or not slides:
        return jsonify({'error': 'Missing required fields'}), 400
    
    # If updating an existing presentation
    if presentation_id:
        presentation = PresentationModel.query.filter_by(id=presentation_id, user_id=user_id).first()
        
        if not presentation:
            return jsonify({'error': 'Presentation not found or unauthorized'}), 404
        
        # Update presentation
        presentation.topic = topic
        presentation.template_id = template_id
        presentation.slide_count = len(slides)
        presentation.updated_at = datetime.utcnow()
        
        # Delete existing slides
        Slide.query.filter_by(presentation_id=presentation.id).delete()
    else:
        # Create new presentation
        presentation = PresentationModel(
            user_id=user_id,
            topic=topic,
            template_id=template_id,
            slide_count=len(slides)
        )
        db.session.add(presentation)
        db.session.flush()  # To get the presentation ID
    
    # Add/update slides
    for i, slide_data in enumerate(slides):
        slide = Slide(
            presentation_id=presentation.id,
            slide_order=i,
            layout=slide_data.get('layout'),
            content=slide_data.get('content', {})
        )
        db.session.add(slide)
    
    db.session.commit()
    
    return jsonify({
        'message': 'Presentation saved successfully',
        'presentation': presentation.to_dict()
    })

@app.route('/api/presentations', methods=['GET'])
@login_required
def list_presentations():
    user_id = session.get('user_id')
    
    # Get all presentations for the user
    presentations = PresentationModel.query.filter_by(user_id=user_id).order_by(PresentationModel.updated_at.desc()).all()
    
    return jsonify({
        'presentations': [p.to_dict() for p in presentations]
    })

@app.route('/api/presentations/<int:presentation_id>', methods=['GET'])
@login_required
def get_presentation(presentation_id):
    user_id = session.get('user_id')
    
    # Get the presentation
    presentation = PresentationModel.query.filter_by(id=presentation_id, user_id=user_id).first()
    
    if not presentation:
        return jsonify({'error': 'Presentation not found'}), 404
    
    return jsonify({
        'presentation': presentation.to_dict()
    })

@app.route('/api/presentations/<int:presentation_id>', methods=['DELETE'])
@login_required
def delete_presentation(presentation_id):
    user_id = session.get('user_id')
    
    # Get the presentation
    presentation = PresentationModel.query.filter_by(id=presentation_id, user_id=user_id).first()
    
    if not presentation:
        return jsonify({'error': 'Presentation not found'}), 404
    
    # Delete the presentation
    db.session.delete(presentation)
    db.session.commit()
    
    return jsonify({
        'message': 'Presentation deleted successfully'
    })

@app.route('/api/export', methods=['POST'])
def export_pptx():
    data = request.json
    slides = data.get('slides', [])
    template_id = data.get('template')
    
    # Create a temporary file
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
        temp_filename = temp_file.name
    
    # Generate the PPTX file
    create_presentation(temp_filename, slides, template_id)
    
    # Send the file
    return send_file(
        temp_filename,
        as_attachment=True,
        download_name=f"{data.get('topic', 'presentation')}.pptx",
        mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation'
    )

def hex_to_rgb(hex_color):
    """Convert hex color to RGB tuple"""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def apply_template_to_slide(slide, template_id):
    """Apply template styling to a slide"""
    template = TEMPLATES.get(template_id, TEMPLATES['corporate'])
    
    # Set background color
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = template['colors']['background']
    
    return template

def apply_text_formatting(paragraph, font_size, color, font_name, bold=False, italic=False, alignment=PP_ALIGN.LEFT):
    """Apply consistent text formatting to a paragraph"""
    paragraph.font.size = Pt(font_size)
    paragraph.font.color.rgb = color
    paragraph.font.name = font_name
    paragraph.font.bold = bold
    paragraph.font.italic = italic
    paragraph.alignment = alignment


def create_title_only_slide(presentation, content, template):
    """Create a title-only slide"""
    slide_layout = presentation.slide_layouts[0]  # Title Slide layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Set title - ensure it's a string
    title = slide.shapes.title
    title_text = content.get('title', 'Title')
    # Check if it's a dictionary (handle special formatting from frontend)
    if isinstance(title_text, dict) and 'text' in title_text:
        title_text = title_text['text']
    # Ensure it's a string
    title_text = str(title_text)
    title.text = title_text
    
    apply_text_formatting(
        title.text_frame.paragraphs[0], 
        font_size=54,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.CENTER
    )
    
    # Set subtitle - ensure it's a string
    subtitle = slide.placeholders[1]
    subtitle_text = content.get('subtitle', 'Subtitle')
    # Check if it's a dictionary
    if isinstance(subtitle_text, dict) and 'text' in subtitle_text:
        subtitle_text = subtitle_text['text']
    # Ensure it's a string
    subtitle_text = str(subtitle_text)
    subtitle.text = subtitle_text
    
    apply_text_formatting(
        subtitle.text_frame.paragraphs[0], 
        font_size=32,
        color=template['colors']['secondary'],
        font_name=template['font'],
        alignment=PP_ALIGN.CENTER
    )
    
    return slide

def create_two_column_slide(presentation, content, template):
    """Create a slide with two columns"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Add title - ensure it's a string
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    
    title_text = content.get('title', 'Title')
    # Check if it's a dictionary
    if isinstance(title_text, dict) and 'text' in title_text:
        title_text = title_text['text']
    # Ensure it's a string
    title_text = str(title_text)
    
    p = text_frame.add_paragraph()
    p.text = title_text
    apply_text_formatting(
        p, 
        font_size=40,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # Column 1 title
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(0.75)
    
    col1_title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = col1_title_box.text_frame
    
    col1_title_text = content.get('column1Title', 'Column 1')
    # Check if it's a dictionary
    if isinstance(col1_title_text, dict) and 'text' in col1_title_text:
        col1_title_text = col1_title_text['text']
    # Ensure it's a string
    col1_title_text = str(col1_title_text)
    
    p = text_frame.add_paragraph()
    p.text = col1_title_text
    apply_text_formatting(
        p, 
        font_size=28,
        color=template['colors']['secondary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # Column 1 content
    left = Inches(0.5)
    top = Inches(2.25)
    width = Inches(4.5)
    height = Inches(3)
    
    col1_content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = col1_content_box.text_frame
    text_frame.word_wrap = True
    
    column1_text = content.get('column1Content', 'Column 1 content')
    # Handle if it's an array, dictionary, or string
    if isinstance(column1_text, list):
        column1_text = "\n".join([str(item) for item in column1_text])
    elif isinstance(column1_text, dict) and 'text' in column1_text:
        column1_text = column1_text['text']
    # Ensure it's a string
    column1_text = str(column1_text)
    
    wrapped_column1 = textwrap.fill(column1_text, width=50)
    
    p = text_frame.add_paragraph()
    p.text = wrapped_column1
    apply_text_formatting(
        p, 
        font_size=20,
        color=template['colors']['text'],
        font_name=template['font'],
        alignment=PP_ALIGN.LEFT
    )
    
    # Column 2 title
    left = Inches(5.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(0.75)
    
    col2_title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = col2_title_box.text_frame
    
    col2_title_text = content.get('column2Title', 'Column 2')
    # Check if it's a dictionary
    if isinstance(col2_title_text, dict) and 'text' in col2_title_text:
        col2_title_text = col2_title_text['text']
    # Ensure it's a string
    col2_title_text = str(col2_title_text)
    
    p = text_frame.add_paragraph()
    p.text = col2_title_text
    apply_text_formatting(
        p, 
        font_size=28,
        color=template['colors']['secondary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # Column 2 content
    left = Inches(5.5)
    top = Inches(2.25)
    width = Inches(4.5)
    height = Inches(3)
    
    col2_content_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = col2_content_box.text_frame
    text_frame.word_wrap = True
    
    column2_text = content.get('column2Content', 'Column 2 content')
    # Handle if it's an array, dictionary, or string
    if isinstance(column2_text, list):
        column2_text = "\n".join([str(item) for item in column2_text])
    elif isinstance(column2_text, dict) and 'text' in column2_text:
        column2_text = column2_text['text']
    # Ensure it's a string
    column2_text = str(column2_text)
    
    wrapped_column2 = textwrap.fill(column2_text, width=50)
    
    p = text_frame.add_paragraph()
    p.text = wrapped_column2
    apply_text_formatting(
        p, 
        font_size=20,
        color=template['colors']['text'],
        font_name=template['font'],
        alignment=PP_ALIGN.LEFT
    )
    
    return slide

def create_image_and_paragraph_slide(presentation, content, template):
    """Create a slide with an image placeholder and paragraph"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Add title - ensure it's a string
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9)
    height = Inches(1)
    
    title_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = title_box.text_frame
    
    title_text = content.get('title', 'Title')
    # Check if it's a dictionary
    if isinstance(title_text, dict) and 'text' in title_text:
        title_text = title_text['text']
    # Ensure it's a string
    title_text = str(title_text)
    
    p = text_frame.add_paragraph()
    p.text = title_text
    apply_text_formatting(
        p, 
        font_size=40,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # Add paragraph
    left = Inches(0.5)
    top = Inches(1.5)
    width = Inches(4.5)
    height = Inches(3.5)
    
    text_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = text_box.text_frame
    text_frame.word_wrap = True
    
    paragraph_text = content.get('paragraph', 'Paragraph text')
    # Check if it's a dictionary or list
    if isinstance(paragraph_text, dict) and 'text' in paragraph_text:
        paragraph_text = paragraph_text['text']
    elif isinstance(paragraph_text, list):
        paragraph_text = "\n".join([str(item) for item in paragraph_text])
    # Ensure it's a string
    paragraph_text = str(paragraph_text)
    
    wrapped_text = textwrap.fill(paragraph_text, width=60)
    
    p = text_frame.add_paragraph()
    p.text = wrapped_text
    apply_text_formatting(
        p, 
        font_size=20,
        color=template['colors']['text'],
        font_name=template['font'],
        alignment=PP_ALIGN.LEFT
    )
    
    # Add image placeholder
    left = Inches(5.5)
    top = Inches(1.5)
    width = Inches(4)
    height = Inches(3.5)
    
    img_placeholder = slide.shapes.add_shape(
        1,  # Rectangle
        left, top, width, height
    )
    img_placeholder.fill.solid()
    img_placeholder.fill.fore_color.rgb = template['colors']['secondary']
    img_placeholder.line.color.rgb = template['colors']['primary']
    
    # Add image description
    left = Inches(5.5)
    top = Inches(3)
    width = Inches(4)
    height = Inches(0.75)
    
    desc_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = desc_box.text_frame
    text_frame.word_wrap = True
    
    img_desc = content.get('imageDescription', 'Image description')
    # Check if it's a dictionary
    if isinstance(img_desc, dict) and 'text' in img_desc:
        img_desc = img_desc['text']
    # Ensure it's a string
    img_desc = str(img_desc)
    
    p = text_frame.add_paragraph()
    p.text = img_desc
    apply_text_formatting(
        p, 
        font_size=16,
        color=template['colors']['background'],
        font_name=template['font'],
        alignment=PP_ALIGN.CENTER
    )
    
    return slide

def create_quote_slide(presentation, content, template):
    """Create a slide with a quote"""
    slide_layout = presentation.slide_layouts[6]  # Blank layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Add quote text
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)
    
    quote_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = quote_box.text_frame
    text_frame.word_wrap = True
    
    # Split the quote into shorter lines if needed
    quote_text = content.get("quote", "Quote goes here")
    # Check if it's a dictionary
    if isinstance(quote_text, dict) and 'text' in quote_text:
        quote_text = quote_text['text']
    # Ensure it's a string
    quote_text = str(quote_text)
    
    wrapped_quote = textwrap.fill(quote_text, width=70)
    
    p = text_frame.add_paragraph()
    p.text = f'"{wrapped_quote}"'
    apply_text_formatting(
        p, 
        font_size=32,
        color=template['colors']['primary'],
        font_name=template['font'],
        italic=True,
        alignment=PP_ALIGN.CENTER
    )
    
    # Add author
    left = Inches(5)
    top = Inches(4.5)
    width = Inches(4)
    height = Inches(1)
    
    author_box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = author_box.text_frame
    
    author_text = content.get("author", "Author")
    # Check if it's a dictionary
    if isinstance(author_text, dict) and 'text' in author_text:
        author_text = author_text['text']
    # Ensure it's a string
    author_text = str(author_text)
    
    p = text_frame.add_paragraph()
    p.text = f'— {author_text}'
    apply_text_formatting(
        p, 
        font_size=24,
        color=template['colors']['secondary'],
        font_name=template['font'],
        alignment=PP_ALIGN.RIGHT
    )
    
    return slide

def create_title_and_bullets_slide(presentation, content, template):
    """Create a slide with title and bullet points"""
    slide_layout = presentation.slide_layouts[1]  # Title and Content layout
    slide = presentation.slides.add_slide(slide_layout)
    
    # Apply template
    template = apply_template_to_slide(slide, template)
    
    # Set title - ensure it's a string
    title = slide.shapes.title
    title_text = content.get('title', 'Title')
    # Check if it's a dictionary
    if isinstance(title_text, dict) and 'text' in title_text:
        title_text = title_text['text']
    # Ensure it's a string
    title_text = str(title_text)
    title.text = title_text
    
    apply_text_formatting(
        title.text_frame.paragraphs[0], 
        font_size=40,
        color=template['colors']['primary'],
        font_name=template['font'],
        bold=True,
        alignment=PP_ALIGN.LEFT
    )
    
    # Set bullet points
    content_shape = slide.placeholders[1]
    text_frame = content_shape.text_frame
    text_frame.clear()
    text_frame.word_wrap = True
    
    # Ensure proper vertical alignment
    text_frame.vertical_anchor = MSO_ANCHOR.TOP
    
    bullets = content.get('bullets', [])
    # Process bullets to ensure they are strings
    processed_bullets = []
    
    if isinstance(bullets, list):
        for bullet in bullets:
            if isinstance(bullet, dict) and 'text' in bullet:
                processed_bullets.append(bullet['text'])
            else:
                processed_bullets.append(str(bullet))
    elif isinstance(bullets, str):
        processed_bullets = [bullets]
    else:
        processed_bullets = ["No bullet points available"]
    
    for bullet in processed_bullets:
        p = text_frame.add_paragraph()
        p.text = bullet
        apply_text_formatting(
            p, 
            font_size=24,
            color=template['colors']['text'],
            font_name=template['font'],
            alignment=PP_ALIGN.LEFT
        )
        p.level = 0
    
    return slide

def create_presentation(filename, slides, template_id):
    """Create a PowerPoint presentation with multiple slides"""
    prs = Presentation()
    
    # Set the slide size to 16:9 aspect ratio
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    # Create slides based on their layout type
    for slide_data in slides:
        layout = slide_data.get('layout')
        content = slide_data.get('content', {})
        
        if layout == 'titleAndBullets':
            create_title_and_bullets_slide(prs, content, template_id)
        elif layout == 'quote':
            create_quote_slide(prs, content, template_id)
        elif layout == 'imageAndParagraph':
            create_image_and_paragraph_slide(prs, content, template_id)
        elif layout == 'twoColumn':
            create_two_column_slide(prs, content, template_id)
        elif layout == 'titleOnly':
            create_title_only_slide(prs, content, template_id)
    
    # Save the presentation
    prs.save(filename)
    
    return filename

# Create presentation routes for handling specific endpoints
@app.route('/presentations')
@login_required
def presentations_list():
    user_id = session.get('user_id')
    username = session.get('username')
    
    # Get all presentations for the user
    presentations = PresentationModel.query.filter_by(user_id=user_id).order_by(PresentationModel.updated_at.desc()).all()
    
    return render_template('dashboard.html', 
                          username=username,
                          presentations=presentations)

# Initialize database tables
with app.app_context():
    db.create_all()

if __name__ == '__main__':
    app.run(debug=True)